from pptx import *
prs = Presentation('template.pptx')
slide = prs.slides[0]
for paragraph in shape.text_frame.paragraphs:
    for run in paragraph.runs:
        print(run.text)
